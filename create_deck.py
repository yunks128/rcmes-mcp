"""Generate RCMES-MCP slide deck using the NASA JPL external-facing template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy
from lxml import etree

TEMPLATE = "/home/ks/science-model-dashboard/NASAjpl_Template_16x9_vA9.pptx"
OUTPUT = "/home/ks/rcmes-mcp/RCMES-MCP_Introduction.pptx"

prs = Presentation(TEMPLATE)

# Remove existing slides (iterate in reverse)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[-1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        # try without namespace
        rId_elem = prs.slides._sldIdLst[-1]
        for rel in prs.part.rels.values():
            if rel.target_part == prs.slides[-1].part:
                rId = rel.rId
                break
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

# Layout references
COVER_NO_IMAGE = prs.slide_layouts[0]    # Cover_No Image
COVER_SQ_IMAGE = prs.slide_layouts[1]    # Cover_Sq Image
TITLE_CONTENT = prs.slide_layouts[5]     # Title and Content
TITLE_SUBTITLE = prs.slide_layouts[6]    # Title and Subtitle
TITLE_SUB_CONTENT = prs.slide_layouts[7] # Title, Subtitle and Content
TWO_CONTENT = prs.slide_layouts[8]       # Two Content
TITLE_ONLY = prs.slide_layouts[4]        # Title Only
BLANK = prs.slide_layouts[17]            # Blank
CLOSING = prs.slide_layouts[18]          # Closing Slide

NAVY = RGBColor(0x00, 0x32, 0x5B)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
FONT = "Arial"

FOOTER_TEXT = "For required markings, please visit https://mh.jpl.nasa.gov"


def set_placeholder_text(slide, idx, text):
    """Set text in a placeholder by index."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return ph
    return None


def add_textbox(slide, left, top, width, height, text, font_size=Pt(14),
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name=FONT):
    """Add a simple textbox."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return box


def set_footer(slide):
    """Set footer placeholder if it exists."""
    for ph in slide.placeholders:
        pf = ph.placeholder_format
        if pf.type is not None and pf.type.name == 'FOOTER':
            ph.text = FOOTER_TEXT


def add_bullet_content(tf, items, font_size=Pt(14), color=DARK_GRAY, spacing=Pt(6)):
    """Add bullet items to an existing text frame."""
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = spacing
        p.level = 0


# ============================================================
# SLIDE 1: Cover — Cover_No Image layout
# ph[15]=section label, ph[16]=title, ph[17]=subtitle, ph[18]=author, ph[20]=footer
# ============================================================
slide = prs.slides.add_slide(COVER_NO_IMAGE)
set_placeholder_text(slide, 15, "Project Update")
set_placeholder_text(slide, 16, "RCMES-MCP")
set_placeholder_text(slide, 17, "Enabling AI-Driven Climate Analysis\nwith NASA's NEX-GDDP-CMIP6 Data")
set_placeholder_text(slide, 18, "March 2026")
set_placeholder_text(slide, 20, FOOTER_TEXT)


# ============================================================
# SLIDE 2: What is RCMES-MCP? — Title and Content
# ph[17]=section, ph[0]=title, ph[19]=content, ph[21]=footer
# ============================================================
slide = prs.slides.add_slide(TITLE_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "What is RCMES-MCP?")
set_footer(slide)

ph = None
for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        ph = p
        break
tf = ph.text_frame
add_bullet_content(tf, [
    "An MCP (Model Context Protocol) server that exposes NASA's climate analysis tools to AI agents",
    "Provides direct access to the NEX-GDDP-CMIP6 dataset — 38 TB of downscaled global climate projections at 0.25° resolution",
    "Enables conversational climate analysis: users ask questions in natural language, AI executes the science",
    "Built on NASA JPL's Regional Climate Model Evaluation System (RCMES) heritage",
    "No data downloads, no coding required — AI handles data access, processing, analysis, and visualization",
], font_size=Pt(14))


# ============================================================
# SLIDE 3: The Data — Two Content layout
# ph[17]=section, ph[0]=title, ph[14]=subtitle, ph[19]=left content, ph[20]=right content
# ============================================================
slide = prs.slides.add_slide(TWO_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "The Data: NEX-GDDP-CMIP6")
set_placeholder_text(slide, 14, "38 TB of bias-corrected, downscaled CMIP6 climate projections hosted on AWS Open Data")
set_footer(slide)

# Left content - dataset details
for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        add_bullet_content(tf, [
            "0.25° resolution (~25 km), daily timestep",
            "Global coverage, 1950–2100",
            "35 CMIP6 climate models",
            "Historical + 4 SSP scenarios",
            "Accessed via S3 in Zarr format — no downloads",
        ], font_size=Pt(13))

# Right content - variables
for p in slide.placeholders:
    if p.placeholder_format.idx == 20:
        tf = p.text_frame
        items = [
            "tas / tasmax / tasmin — Temperature",
            "pr — Precipitation",
            "hurs / huss — Humidity",
            "sfcWind — Surface wind speed",
            "rsds / rlds — Radiation (SW & LW)",
        ]
        # Add header
        p0 = tf.paragraphs[0]
        p0.text = "9 Climate Variables:"
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.font.name = FONT
        p0.space_after = Pt(4)
        for item in items:
            pp = tf.add_paragraph()
            pp.text = item
            pp.font.size = Pt(12)
            pp.font.color.rgb = DARK_GRAY
            pp.font.name = FONT
            pp.space_after = Pt(3)


# ============================================================
# SLIDE 4: Capabilities — Title and Content
# ============================================================
slide = prs.slides.add_slide(TITLE_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "32+ Climate Analysis Tools")
set_footer(slide)

for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        items = [
            "Data Access — Load climate data by model, variable, scenario; region/country selection; dataset listing",
            "Processing — Temporal & spatial subsetting, regridding, unit conversion, anomaly calculation, country masking",
            "Statistical Analysis — Climatology, trend analysis with significance testing, bias, RMSE, correlation",
            "Extreme Indices — 22 ETCCDI indices (temperature & precipitation extremes), heatwave & drought analysis",
            "Model Evaluation — Multi-model comparison, Taylor diagrams, bias maps",
            "Visualization — Spatial maps, time series, comparison plots, histograms, country-specific maps",
        ]
        add_bullet_content(tf, items, font_size=Pt(13))


# ============================================================
# SLIDE 5: How It Works — Title and Content
# ============================================================
slide = prs.slides.add_slide(TITLE_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "How It Works")
set_footer(slide)

for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        items = [
            "User asks a climate question in natural language (e.g., \"What is the heatwave trend in California?\")",
            "AI agent selects and chains appropriate tools via the MCP protocol",
            "Tools access NEX-GDDP-CMIP6 data directly from AWS S3 — no downloads needed",
            "Analysis engine (xarray, Dask, xclim, SciPy) processes the data",
            "Results returned as statistics, maps, time series plots, or derived datasets",
            "Chainable operations: each tool returns a dataset_id for multi-step workflows",
        ]
        add_bullet_content(tf, items, font_size=Pt(13))

# Add flow diagram below using textboxes
# User → MCP Server → AWS S3 → Analysis → Results
flow_labels = ["User / AI Agent", "MCP Server", "AWS S3 Data", "Analysis Engine", "Results"]
flow_colors = [
    RGBColor(0x00, 0x7D, 0xBA),  # blue
    NAVY,
    RGBColor(0x23, 0x72, 0x3B),  # green
    RGBColor(0xE6, 0x7E, 0x22),  # orange
    RGBColor(0xC0, 0x39, 0x2B),  # red
]
box_w = Inches(1.65)
box_h = Inches(0.45)
gap = Inches(0.18)
total = len(flow_labels) * box_w + (len(flow_labels) - 1) * gap
start_x = (prs.slide_width - total) // 2
y = Inches(4.55)

for i, (label, color) in enumerate(zip(flow_labels, flow_colors)):
    x = start_x + i * (box_w + gap)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    stf = shape.text_frame
    stf.vertical_anchor = MSO_ANCHOR.MIDDLE
    sp = stf.paragraphs[0]
    sp.text = label
    sp.font.size = Pt(9)
    sp.font.bold = True
    sp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sp.font.name = FONT
    sp.alignment = PP_ALIGN.CENTER

    if i < len(flow_labels) - 1:
        ax = x + box_w + Emu(10000)
        ay = y + box_h // 2 - Inches(0.08)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, gap - Emu(20000), Inches(0.16))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
        arrow.line.fill.background()


# ============================================================
# SLIDE 6: Example Use Case — Title, Subtitle and Content
# ============================================================
slide = prs.slides.add_slide(TITLE_SUB_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "Example: Heatwave Trend Analysis")
set_placeholder_text(slide, 14, "\"What is the heatwave trend in California under the high-emissions scenario?\"")
set_footer(slide)

for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        steps = [
            "1. AI loads data: load_climate_data(tasmax, ACCESS-CM2, ssp585, California, 2015–2100)",
            "2. AI analyzes heatwaves: analyze_heatwaves(dataset_id, threshold_percentile=90)",
            "3. AI computes trend: calculate_trend(heatwave_id) → significant warming trend with p-value",
            "4. AI visualizes: generate_timeseries_plot(dataset_id, show_trend=True) → plot returned to user",
            "",
            "All steps executed automatically — the user only asks the question",
        ]
        add_bullet_content(tf, steps, font_size=Pt(13))
        # Make last item italic
        last_p = tf.paragraphs[-1]
        last_p.font.italic = True
        last_p.font.color.rgb = RGBColor(0x00, 0x7D, 0xBA)


# ============================================================
# SLIDE 7: Current Status — Two Content
# ============================================================
slide = prs.slides.add_slide(TWO_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "Current Status")
set_placeholder_text(slide, 14, "")
set_footer(slide)

# Left - accomplishments
for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        p0 = tf.paragraphs[0]
        p0.text = "Accomplishments"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.font.name = FONT
        for item in [
            "32+ MCP tools: full climate analysis pipeline",
            "Operational 38 TB data access via AWS S3",
            "22 ETCCDI extreme climate indices",
            "Interactive web UI (React + FastAPI)",
            "Multiple deployment modes (MCP, Web, API, Docker)",
            "Session management with chainable operations",
            "Country-level masking & geographic selection",
        ]:
            pp = tf.add_paragraph()
            pp.text = item
            pp.font.size = Pt(11)
            pp.font.color.rgb = DARK_GRAY
            pp.font.name = FONT
            pp.space_after = Pt(3)

# Right - tech stack
for p in slide.placeholders:
    if p.placeholder_format.idx == 20:
        tf = p.text_frame
        p0 = tf.paragraphs[0]
        p0.text = "Technology Stack"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.font.name = FONT
        stack = [
            ("Core:", "Python 3.10+, xarray, Dask, zarr"),
            ("Climate:", "xclim, SciPy, NetCDF4"),
            ("Geospatial:", "Shapely, GeoPandas, Cartopy"),
            ("Server:", "FastMCP, FastAPI, Uvicorn"),
            ("Frontend:", "React, TypeScript"),
            ("Data:", "AWS S3, NEX-GDDP-CMIP6"),
            ("AI:", "MCP Protocol, Azure OpenAI"),
        ]
        for label, tech in stack:
            pp = tf.add_paragraph()
            run1 = pp.add_run()
            run1.text = f"{label} "
            run1.font.size = Pt(11)
            run1.font.bold = True
            run1.font.color.rgb = RGBColor(0x00, 0x7D, 0xBA)
            run1.font.name = FONT
            run2 = pp.add_run()
            run2.text = tech
            run2.font.size = Pt(11)
            run2.font.color.rgb = DARK_GRAY
            run2.font.name = FONT
            pp.space_after = Pt(3)


# ============================================================
# SLIDE 8: Path Forward — Two Content
# ============================================================
slide = prs.slides.add_slide(TWO_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "Path Forward & Continued Support")
set_placeholder_text(slide, 14, "")
set_footer(slide)

# Left - growth opportunities
for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        p0 = tf.paragraphs[0]
        p0.text = "Growth Opportunities"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.font.name = FONT
        for item in [
            "Expand to additional NASA datasets (MERRA-2, ECCO, GLDAS)",
            "Model evaluation with observational datasets",
            "NASA Earthdata integration for authenticated access",
            "Multi-model ensemble analysis & weighting",
            "Community tool for broader climate science use",
            "Reproducible climate assessments for policy support",
        ]:
            pp = tf.add_paragraph()
            pp.text = item
            pp.font.size = Pt(11)
            pp.font.color.rgb = DARK_GRAY
            pp.font.name = FONT
            pp.space_after = Pt(3)

# Right - impact
for p in slide.placeholders:
    if p.placeholder_format.idx == 20:
        tf = p.text_frame
        p0 = tf.paragraphs[0]
        p0.text = "What Continued Support Enables"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.font.name = FONT
        for item in [
            "Democratize access to NASA climate data for non-experts",
            "Accelerate climate research by removing data engineering bottlenecks",
            "Enable rapid AI-assisted climate impact assessments",
            "Lower the barrier to entry across disciplines",
            "Position NASA JPL at the forefront of AI-enabled science",
            "Support decision-makers with actionable climate intelligence",
        ]:
            pp = tf.add_paragraph()
            pp.text = item
            pp.font.size = Pt(11)
            pp.font.color.rgb = DARK_GRAY
            pp.font.name = FONT
            pp.space_after = Pt(3)


# ============================================================
# SLIDE 9: Closing slide
# ============================================================
slide = prs.slides.add_slide(CLOSING)
set_footer(slide)


# ============================================================
# Save
# ============================================================
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
