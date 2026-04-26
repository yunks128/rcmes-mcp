"""Generate RCMES-MCP slide deck using the NASA JPL external-facing template.

Also regenerates illustrative PNG figures (under docs/figures/) for the advanced
spatiotemporal-analysis slides. Figures use synthetic data — they document what
each tool produces, not real model output.
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy
from lxml import etree

# Headless matplotlib for figure generation
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import numpy as np

TEMPLATE = "/home/ks/science-model-dashboard/NASAjpl_Template_16x9_vA9.pptx"
OUTPUT = "/home/ks/rcmes-mcp/docs/RCMES-MCP_Introduction.pptx"
FIG_DIR = Path("/home/ks/rcmes-mcp/docs/figures")
FIG_DIR.mkdir(parents=True, exist_ok=True)

prs = Presentation(TEMPLATE)

# Remove existing slides (iterate in reverse)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[-1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        # try without namespace
        rId_elem = prs.slides._sldIdLst[-1]
        for rel in prs.part.rels.values():
            if rel.target_part == prs.slides[-1].part:
                rId = rel.rId
                break
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

# Layout references
COVER_NO_IMAGE = prs.slide_layouts[0]    # Cover_No Image
COVER_SQ_IMAGE = prs.slide_layouts[1]    # Cover_Sq Image
TITLE_CONTENT = prs.slide_layouts[5]     # Title and Content
TITLE_SUBTITLE = prs.slide_layouts[6]    # Title and Subtitle
TITLE_SUB_CONTENT = prs.slide_layouts[7] # Title, Subtitle and Content
TWO_CONTENT = prs.slide_layouts[8]       # Two Content
TITLE_ONLY = prs.slide_layouts[4]        # Title Only
BLANK = prs.slide_layouts[17]            # Blank
CLOSING = prs.slide_layouts[18]          # Closing Slide

NAVY = RGBColor(0x00, 0x32, 0x5B)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
FONT = "Arial"

FOOTER_TEXT = "For required markings, please visit https://mh.jpl.nasa.gov"


def set_placeholder_text(slide, idx, text):
    """Set text in a placeholder by index."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return ph
    return None


def add_textbox(slide, left, top, width, height, text, font_size=Pt(14),
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name=FONT):
    """Add a simple textbox."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return box


def set_footer(slide):
    """Set footer placeholder if it exists."""
    for ph in slide.placeholders:
        pf = ph.placeholder_format
        if pf.type is not None and pf.type.name == 'FOOTER':
            ph.text = FOOTER_TEXT


def add_bullet_content(tf, items, font_size=Pt(14), color=DARK_GRAY, spacing=Pt(6)):
    """Add bullet items to an existing text frame."""
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = spacing
        p.level = 0


# ============================================================
# FIGURE GENERATORS — synthetic but representative outputs of each new tool
# Saved to docs/figures/, embedded into the new slides below.
# ============================================================

_RNG = np.random.default_rng(42)


def _save(fig, name):
    path = FIG_DIR / name
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return str(path)


# ----------------------------------------------------------------------
# Real-data path: load NEX-GDDP-CMIP6 via RCMES tools, run actual analyses.
# Falls back to synthetic figures on any failure so deck regen never blocks.
# Toggle with RCMES_DECK_REAL_DATA=0 to force synthetic.
# ----------------------------------------------------------------------
os.environ.setdefault("RCMES_CACHE_DIR", "/home/ks/rcmes-mcp/rcmes_mcp_cache")
USE_REAL_DATA = os.environ.get("RCMES_DECK_REAL_DATA", "1") != "0"
_REAL_AVAILABLE = False
if USE_REAL_DATA:
    try:
        import xarray as xr
        from scipy import ndimage
        from rcmes_mcp.tools.data_access import load_climate_data, wait_for_materialization
        from rcmes_mcp.utils.session import session_manager
        _REAL_AVAILABLE = True
        print("[deck] Real-data mode ON — figures will use NEX-GDDP-CMIP6 from S3")
    except Exception as e:
        print(f"[deck] Real-data deps unavailable ({e}); using synthetic figures")
else:
    print("[deck] Synthetic figures requested (RCMES_DECK_REAL_DATA=0)")

CALIFORNIA = dict(lat_min=32.0, lat_max=42.0, lon_min=-125.0, lon_max=-114.0)
LA_BASIN = dict(lat_min=32.0, lat_max=36.0, lon_min=-120.0, lon_max=-116.0)
WESTERN_US = dict(lat_min=30.0, lat_max=50.0, lon_min=-125.0, lon_max=-110.0)


def _load(variable, model, scenario, start_date, end_date, **bbox):
    """Load NEX-GDDP-CMIP6 via the RCMES load_climate_data tool. Returns DataArray."""
    print(f"[deck] _load {variable} {model}/{scenario} {start_date}..{end_date} bbox={bbox}")
    r = load_climate_data(
        variable=variable, model=model, scenario=scenario,
        start_date=start_date, end_date=end_date, **bbox,
    )
    if not r.get("success"):
        raise RuntimeError(f"load_climate_data failed: {r.get('error')}")
    ds_id = r["dataset_id"]
    wait_for_materialization(ds_id, timeout=1800)
    ds = session_manager.get(ds_id)
    if isinstance(ds, xr.DataArray):
        return ds
    var = variable if variable in ds.data_vars else list(ds.data_vars)[0]
    return ds[var]


def _normalize_lons_arr(arr_2d, lons):
    """Convert 0-360 lons to -180:180 and re-sort columns of arr_2d to match."""
    lons = np.asarray(lons)
    if lons.max() > 180:
        lons = ((lons + 180) % 360) - 180
        sort = np.argsort(lons)
        return arr_2d[..., sort], lons[sort]
    return arr_2d, lons


def _to_year_index(da):
    """Replace time coord with integer year so xr.concat works across CMIP6 calendars."""
    years = np.array([
        t.year if hasattr(t, "year") else int(str(t)[:4])
        for t in da.time.values
    ])
    return da.assign_coords(time=years)


# ============================================================
# REAL-DATA FIGURE GENERATORS
# ============================================================

def _real_anomaly_hovmoller():
    # Near-future window so individual extreme blobs stand out against
    # a near-zero background rather than uniformly-red late-century forcing.
    fut = _load("tasmax", "ACCESS-CM2", "ssp585",
                "2025-01-01", "2029-12-31", **CALIFORNIA)
    base = _load("tasmax", "ACCESS-CM2", "historical",
                 "1980-01-01", "2009-12-31", **CALIFORNIA)
    print("[deck] Computing day-of-year z-score climatology...")
    clim_mean = base.groupby("time.dayofyear").mean()
    clim_std = base.groupby("time.dayofyear").std().where(lambda s: s > 0)
    z = (fut.groupby("time.dayofyear") - clim_mean).groupby("time.dayofyear") / clim_std
    z = z.resample(time="7D").mean()  # weekly for cleaner Hovmöller
    hov = z.mean(dim="lat").compute()

    arr, lons = _normalize_lons_arr(hov.values, hov.lon.values)
    times = hov.time.values

    # Find blob events on the (time × lon) hovmöller
    mask = np.abs(arr) > 1.5
    labels, n = ndimage.label(mask)
    boxes = []
    for i in range(1, n + 1):
        idx = np.where(labels == i)
        if idx[0].size < 8:
            continue
        boxes.append((int(idx[1].min()), int(idx[1].max()),
                      int(idx[0].min()), int(idx[0].max())))

    fig, ax = plt.subplots(figsize=(8.5, 5.4))
    vlim = 3.5
    im = ax.pcolormesh(lons, times, arr, cmap="RdBu_r",
                       vmin=-vlim, vmax=vlim, shading="auto")
    for x0, x1, t0, t1 in boxes:
        ax.plot(
            [lons[x0], lons[x1], lons[x1], lons[x0], lons[x0]],
            [times[t0], times[t0], times[t1], times[t1], times[t0]],
            color="black", linewidth=1.2, linestyle="--",
        )
    ax.set_xlabel("Longitude (°)")
    ax.set_ylabel("Time")
    ax.set_title("California tasmax z-score — zonal mean Hovmöller\n"
                 "ACCESS-CM2 SSP5-8.5 (2025–29) vs 1980–2009 baseline")
    cb = plt.colorbar(im, ax=ax, pad=0.02)
    cb.set_label("z-score (σ)")
    return _save(fig, "fig_anomaly_hovmoller.png")


def _real_eof():
    base = _load("tasmax", "ACCESS-CM2", "historical",
                 "1980-01-01", "2009-12-31", **CALIFORNIA)
    print("[deck] Monthly anomaly + EOF SVD...")
    monthly = base.resample(time="MS").mean()
    clim = monthly.groupby("time.month").mean()
    anom = (monthly.groupby("time.month") - clim).compute().transpose("time", "lat", "lon")

    arr = anom.values.astype(np.float64)
    n_t, n_lat, n_lon = arr.shape
    lat_w = np.sqrt(np.cos(np.deg2rad(anom.lat.values))).clip(0)
    weighted = arr * lat_w[np.newaxis, :, np.newaxis]
    flat = weighted.reshape(n_t, -1)
    col_finite = np.isfinite(flat).all(axis=0)
    flat_clean = np.nan_to_num(flat[:, col_finite])

    U, S, Vt = np.linalg.svd(flat_clean, full_matrices=False)
    var_frac = (S ** 2) / np.sum(S ** 2)

    eof1_flat = np.full(n_lat * n_lon, np.nan)
    eof1_flat[col_finite] = Vt[0, :]
    eof1 = (eof1_flat.reshape(n_lat, n_lon) / lat_w[:, np.newaxis])
    pc1 = U[:, 0] * S[0]

    eof1, lons = _normalize_lons_arr(eof1, anom.lon.values)
    lats = anom.lat.values

    fig = plt.figure(figsize=(11, 4.5))
    ax_map = fig.add_subplot(1, 2, 1)
    vlim = float(np.nanmax(np.abs(eof1)))
    im = ax_map.pcolormesh(lons, lats, eof1, cmap="RdBu_r",
                           vmin=-vlim, vmax=vlim, shading="auto")
    ax_map.set_xlabel("Longitude (°)")
    ax_map.set_ylabel("Latitude (°)")
    ax_map.set_title(f"EOF1 — California tasmax monthly anomaly ({var_frac[0]*100:.0f}% var)\nACCESS-CM2 historical 1980–2009")
    plt.colorbar(im, ax=ax_map, pad=0.02, label="loading")

    ax_ts = fig.add_subplot(1, 2, 2)
    months = anom.time.values
    ax_ts.plot(months, pc1, color="#0f4c81", linewidth=1.0)
    ax_ts.axhline(0, color="grey", linewidth=0.6)
    ax_ts.fill_between(months, pc1, 0, where=pc1 > 0, color="#c4321c", alpha=0.3)
    ax_ts.fill_between(months, pc1, 0, where=pc1 < 0, color="#1a73c2", alpha=0.3)
    ax_ts.set_xlabel("Time")
    ax_ts.set_ylabel("PC1 amplitude")
    ax_ts.set_title("Principal component 1")
    ax_ts.grid(alpha=0.3)
    fig.tight_layout()
    return _save(fig, "fig_eof.png")


def _real_ensemble_spread():
    models = ["ACCESS-CM2", "GFDL-ESM4", "MRI-ESM2-0"]
    series = []
    used = []
    for m in models:
        try:
            data = _load("tas", m, "ssp585",
                         "2050-01-01", "2099-12-31", **LA_BASIN)
            ts = data.mean(dim=["lat", "lon"]).resample(time="1YS").mean().compute()
            ts = _to_year_index(ts)  # normalize across model calendars
            series.append(ts)
            used.append(m)
        except Exception as e:
            print(f"[deck]   {m} skipped: {e}")
    if not series:
        raise RuntimeError("no ensemble members loaded")
    stack = xr.concat(series, dim=xr.Variable("model", used), join="outer")
    if float(stack.mean()) > 200:
        stack = stack - 273.15
    ref = float(stack.median(dim="model").isel(time=0))
    stack_anom = stack - ref
    p05 = stack_anom.quantile(0.05, dim="model")
    p25 = stack_anom.quantile(0.25, dim="model")
    p50 = stack_anom.median(dim="model")
    p75 = stack_anom.quantile(0.75, dim="model")
    p95 = stack_anom.quantile(0.95, dim="model")
    times = stack_anom.time.values

    fig, ax = plt.subplots(figsize=(9, 5))
    ax.fill_between(times, p05, p95, color="#1f77b4", alpha=0.18, label="5–95% across models")
    ax.fill_between(times, p25, p75, color="#1f77b4", alpha=0.3, label="25–75%")
    for ts in series:
        vals = ts.values
        if float(vals.mean()) > 200:
            vals = vals - 273.15
        ax.plot(ts.time.values, vals - ref, color="grey", linewidth=0.5, alpha=0.5)
    ax.plot(times, p50, color="#1f77b4", linewidth=2, label="Median")
    ax.set_xlabel("Year")
    ax.set_ylabel("Δ Temperature vs 2050 (°C)")
    ax.set_title(f"LA-area tas SSP5-8.5 — {len(used)}-model ensemble spread")
    ax.legend(loc="upper left", frameon=True)
    ax.grid(alpha=0.3)
    return _save(fig, "fig_ensemble_spread.png")


def _real_scenario_fan():
    scenarios = [("ssp126", "#1a9850"),
                 ("ssp245", "#fdae61"),
                 ("ssp585", "#7f0000")]
    fig, ax = plt.subplots(figsize=(9, 5))
    ref = None
    plotted = 0
    for sc, color in scenarios:
        try:
            data = _load("tas", "ACCESS-CM2", sc,
                         "2050-01-01", "2099-12-31", **LA_BASIN)
            ts = data.mean(dim=["lat", "lon"]).resample(time="1YS").mean().compute()
            vals = ts.values
            if float(vals.mean()) > 200:
                vals = vals - 273.15
            if ref is None:
                ref = float(vals[0])
            anom = vals - ref
            # 5-yr rolling mean for smoother visual
            kernel = np.ones(5) / 5
            smoothed = np.convolve(anom, kernel, mode="same")
            ax.plot(ts.time.values, anom, color=color, linewidth=0.6, alpha=0.35)
            ax.plot(ts.time.values, smoothed, color=color, linewidth=2.0, label=sc)
            plotted += 1
        except Exception as e:
            print(f"[deck]   scenario {sc} skipped: {e}")
    if plotted == 0:
        raise RuntimeError("no scenarios plotted")
    ax.set_xlabel("Year")
    ax.set_ylabel("Δ Temperature vs 2050 (°C)")
    ax.set_title("Scenario divergence — ACCESS-CM2 LA-area tas")
    ax.legend(loc="upper left", frameon=True)
    ax.grid(alpha=0.3)
    return _save(fig, "fig_scenario_fan.png")


def _real_time_of_emergence():
    hist = _load("tas", "ACCESS-CM2", "historical",
                 "1980-01-01", "1999-12-31", **WESTERN_US)
    # ssp245 (Middle of the Road) gives a slower, more pedagogical emergence pattern —
    # under ssp585 the signal exceeds even 2σ across the entire Western US by ~2025.
    fut = _load("tas", "ACCESS-CM2", "ssp245",
                "2015-01-01", "2099-12-31", **WESTERN_US)
    print("[deck] Computing time-of-emergence map...")
    ah = hist.resample(time="1YS").mean().compute()
    af = fut.resample(time="1YS").mean().compute()
    mu0 = ah.mean(dim="time")
    sigma0 = ah.std(dim="time")
    # Require full 20-yr window (no edge bias) and use a stricter 2σ threshold so
    # the emergence year actually varies across cells instead of pegging at the
    # first valid year everywhere.
    rolled = af.rolling(time=20, center=True, min_periods=20).mean()
    deviation = np.abs(rolled - mu0)
    threshold = 2.0 * sigma0
    emerged = deviation > threshold
    years_da = xr.DataArray(rolled["time"].dt.year.values, dims=["time"], coords={"time": rolled.time})
    emergence_year = years_da.where(emerged).min(dim="time", skipna=True).compute()

    arr, lons = _normalize_lons_arr(emergence_year.values, emergence_year.lon.values)
    lats = emergence_year.lat.values
    # Replace any non-finite (never-emerged or oddities) with NaN, set explicit vmin/vmax
    arr = np.where(np.isfinite(arr), arr, np.nan)
    finite = arr[np.isfinite(arr)]
    if finite.size:
        vmin = float(np.nanpercentile(finite, 5))
        vmax = float(np.nanpercentile(finite, 95))
        if vmax - vmin < 5:
            vmax = vmin + 10
    else:
        vmin, vmax = 2030, 2090
    fig, ax = plt.subplots(figsize=(10, 5))
    im = ax.pcolormesh(lons, lats, arr, cmap="viridis", vmin=vmin, vmax=vmax, shading="auto")
    ax.set_xlabel("Longitude (°)")
    ax.set_ylabel("Latitude (°)")
    ax.set_title("Western US ToE — ACCESS-CM2 SSP2-4.5 (2015–2099)\n"
                 "year |Δtas| > 2σ baseline (1980–1999)")
    cb = plt.colorbar(im, ax=ax, pad=0.02)
    cb.set_label("Emergence year")
    return _save(fig, "fig_time_of_emergence.png")


def _real_change_and_agreement():
    models = ["ACCESS-CM2", "GFDL-ESM4", "MRI-ESM2-0"]
    changes = []
    used = []
    for m in models:
        try:
            h = _load("tas", m, "historical", "1990-01-01", "1999-12-31", **WESTERN_US)
            f = _load("tas", m, "ssp585", "2090-01-01", "2099-12-31", **WESTERN_US)
            change = (f.mean(dim="time") - h.mean(dim="time")).compute()
            # Drop time coord (already reduced); keep only spatial
            change = change.reset_coords(drop=True) if hasattr(change, "reset_coords") else change
            changes.append(change)
            used.append(m)
        except Exception as e:
            print(f"[deck]   change member {m} skipped: {e}")
    if len(changes) < 2:
        raise RuntimeError("need ≥2 models for agreement")
    stack = xr.concat(changes, dim=xr.Variable("model", used), join="outer")
    ens_mean = stack.mean(dim="model")
    ens_sign = np.sign(ens_mean)
    same_sign = (np.sign(stack) == ens_sign)
    agreement = same_sign.sum(dim="model") / len(used)
    # Mask cells where any model is NaN (ocean / outside coverage)
    valid = stack.notnull().all(dim="model")
    agreement = agreement.where(valid)

    arr_mean, lons = _normalize_lons_arr(ens_mean.values, ens_mean.lon.values)
    arr_agree, _ = _normalize_lons_arr(agreement.values, agreement.lon.values)
    lats = ens_mean.lat.values

    fig, axes = plt.subplots(1, 2, figsize=(12, 4.6))
    vlim = float(np.nanmax(np.abs(arr_mean)))
    im1 = axes[0].pcolormesh(lons, lats, arr_mean, cmap="RdBu_r",
                             vmin=-vlim, vmax=vlim, shading="auto")
    axes[0].set_title(f"Ensemble-mean ΔT — Western US ({len(used)} models)\n"
                      "2090-99 SSP5-8.5 vs 1990-99 historical")
    axes[0].set_xlabel("Longitude")
    axes[0].set_ylabel("Latitude")
    plt.colorbar(im1, ax=axes[0], pad=0.02, label="K")

    im2 = axes[1].pcolormesh(lons, lats, arr_agree, cmap="YlGnBu", vmin=0, vmax=1, shading="auto")
    axes[1].set_title("Model agreement on sign")
    axes[1].set_xlabel("Longitude")
    plt.colorbar(im2, ax=axes[1], pad=0.02, label="fraction agreeing")
    fig.tight_layout()
    return _save(fig, "fig_change_and_agreement.png")


# ============================================================
# Pillar 3 — Ensemble weighting (real data via NCEP reference)
# ============================================================

_WEIGHT_METHOD_COLORS = {
    "equal":        "#666666",
    "skill":        "#1a73c2",
    "independence": "#1a9850",
    "combined":     "#c4321c",
}


def _real_weighting_validation():
    """
    Train weights on 1980–2002 historical, validate on 2003–2014 against NCEP
    reanalysis. Then apply the four trained weight vectors to the cached
    SSP5-8.5 2050–2099 ensemble and compare projections.

    Produces both fig_weighting_validation.png and fig_weighting_comparison.png
    in one pass (cheaper — share the trained-weights state).
    """
    from rcmes_mcp.tools.data_access import load_reference_dataset
    from rcmes_mcp.tools.analysis import (
        validate_ensemble_weighting,
        calculate_model_weights,
        apply_ensemble_weights,
    )

    # 1) Load reference (NCEP regional series, very small)
    ref_r = load_reference_dataset(
        source="ncep-reanalysis-monthly",
        start_date="1980-01-01", end_date="2014-12-31",
        lat_min=32.0, lat_max=36.0, lon_min=-120.0, lon_max=-116.0,
    )
    if not ref_r.get("success"):
        raise RuntimeError(f"reference load failed: {ref_r.get('error')}")
    reference_id = ref_r["dataset_id"]
    print(f"[deck]   reference dataset: {reference_id} dims={ref_r['dimensions']}")

    # 2) Build a HISTORICAL ensemble for the same LA-basin bbox (3 models).
    # Pre-resample to monthly + normalize calendar so xr.concat doesn't trip on
    # CMIP6's mixed calendar conventions (noleap vs standard).
    models = ["ACCESS-CM2", "GFDL-ESM4", "MRI-ESM2-0"]
    series_hist = []
    used = []
    for m in models:
        try:
            data = _load("tas", m, "historical",
                         "1980-01-01", "2014-12-31", **LA_BASIN)
            monthly = data.resample(time="MS").mean()
            try:
                monthly = monthly.convert_calendar("standard", align_on="date", missing=np.nan)
            except Exception:
                pass
            monthly = monthly.assign_coords(model=m).expand_dims("model")
            series_hist.append(monthly)
            used.append(m)
        except Exception as e:
            print(f"[deck]   {m} historical load failed: {e}")
    if len(series_hist) < 2:
        raise RuntimeError("need ≥2 historical models")
    hist_ensemble = xr.concat(series_hist, dim="model", join="outer")
    hist_id = session_manager.store(
        data=hist_ensemble, source="NEX-GDDP-CMIP6", variable="tas",
        model=",".join(used), scenario="historical",
        description=f"Hist ensemble {used} 1980-2014 LA basin",
    )

    # 3) Validate
    val = validate_ensemble_weighting(
        ensemble_dataset_id=hist_id,
        reference_dataset_id=reference_id,
        train_start="1980-01-01", train_end="2002-12-31",
        test_start="2003-01-01", test_end="2014-12-31",
        sigma_d=0.5, sigma_s=0.5,
    )
    if not val.get("success"):
        raise RuntimeError(f"validation failed: {val.get('error')}")
    print(f"[deck]   validation per-method:")
    for m, r in val["per_method"].items():
        print(f"     {m:>13s}  weights={r.get('weights')} test_rmse={r.get('test_rmse_K')} corr={r.get('test_correlation')}")

    # ---------- Figure A: validation skill bar chart -----------------------
    methods = list(val["per_method"].keys())
    test_rmse = [val["per_method"][m]["test_rmse_K"] for m in methods]
    test_corr = [val["per_method"][m]["test_correlation"] for m in methods]
    train_rmse = [val["per_method"][m]["train_rmse_K"] for m in methods]
    colors = [_WEIGHT_METHOD_COLORS.get(m, "#888") for m in methods]

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))
    x = np.arange(len(methods))
    bw = 0.35
    axes[0].bar(x - bw / 2, train_rmse, bw, color=colors, alpha=0.55, label="train (1980–2002)")
    axes[0].bar(x + bw / 2, test_rmse, bw, color=colors, label="test (2003–2014)")
    axes[0].set_xticks(x)
    axes[0].set_xticklabels(methods, rotation=15)
    axes[0].set_ylabel("Regional-mean RMSE vs NCEP (K)")
    axes[0].set_title("Skill — lower is better")
    axes[0].grid(alpha=0.3, axis="y")
    axes[0].legend()

    axes[1].bar(x, test_corr, color=colors)
    axes[1].set_xticks(x)
    axes[1].set_xticklabels(methods, rotation=15)
    axes[1].set_ylabel("Pearson r (test window)")
    axes[1].set_title("Correlation — higher is better")
    axes[1].set_ylim(0, 1.05)
    axes[1].grid(alpha=0.3, axis="y")

    fig.suptitle(
        f"Ensemble-weighting validation — LA basin tas, train 1980-2002 / test 2003-2014\n"
        f"({len(used)} models, reference: NCEP Reanalysis 1)",
        fontsize=11,
    )
    fig.tight_layout()
    _save(fig, "fig_weighting_validation.png")

    # ---------- Figure B: future projections under each weighting -----------
    # Apply each trained weight set to the cached SSP5-8.5 2050–2099 ensemble.
    # Reuse the ensemble we built for the spread plot.
    fut_models = used
    fut_series = []
    for m in fut_models:
        data = _load("tas", m, "ssp585",
                     "2050-01-01", "2099-12-31", **LA_BASIN)
        # Same calendar normalization as the historical ensemble
        monthly = data.resample(time="MS").mean()
        try:
            monthly = monthly.convert_calendar("standard", align_on="date", missing=np.nan)
        except Exception:
            pass
        monthly = monthly.assign_coords(model=m).expand_dims("model")
        fut_series.append(monthly)
    fut_ensemble = xr.concat(fut_series, dim="model", join="outer")
    fut_id = session_manager.store(
        data=fut_ensemble, source="NEX-GDDP-CMIP6", variable="tas",
        model=",".join(fut_models), scenario="ssp585",
        description=f"Fut ensemble {fut_models} 2050-99 LA basin",
    )

    fig2, ax = plt.subplots(figsize=(10, 5))
    for m in methods:
        w = val["per_method"][m].get("weights")
        if not w:
            continue
        applied = apply_ensemble_weights(ensemble_dataset_id=fut_id, weights=w)
        if not applied.get("success"):
            print(f"[deck]   apply {m} failed: {applied.get('error')}")
            continue
        wds = session_manager.get(applied["dataset_id"])
        var = wds.name if isinstance(wds, xr.DataArray) else list(wds.data_vars)[0]
        wda = wds if isinstance(wds, xr.DataArray) else wds[var]
        ts = wda.mean(dim=["lat", "lon"]).resample(time="1YS").mean().compute()
        vals = ts.values
        if float(vals.mean()) > 200:
            vals = vals - 273.15
        ref0 = float(vals[0])
        anom = vals - ref0
        # Smooth with 5-yr rolling for readability
        kernel = np.ones(5) / 5
        smoothed = np.convolve(anom, kernel, mode="same")
        color = _WEIGHT_METHOD_COLORS.get(m, "#888")
        ax.plot(ts.time.values, anom, color=color, alpha=0.25, linewidth=0.6)
        # Annotate effective number of weighted members (Neff = 1/Σwᵢ²) and max weight
        n_eff = 1.0 / float(np.sum(np.array(w) ** 2))
        wmax = max(w)
        label = f"{m}  (N_eff={n_eff:.1f}, max w={wmax:.2f})"
        ax.plot(ts.time.values, smoothed, color=color, linewidth=2.0, label=label)
    ax.set_xlabel("Year")
    ax.set_ylabel("Δ T vs 2050 (°C)")
    ax.set_title("Weighted SSP5-8.5 projection — LA basin tas\n"
                 "Each colour applies a different weighting scheme to the same 3-model ensemble")
    ax.grid(alpha=0.3)
    ax.legend(loc="upper left", fontsize=8)
    return _save(fig2, "fig_weighting_comparison.png")


def _real_scenario_weighted_combine():
    """
    Probability-weighted scenario combining: load all 4 SSPs for ACCESS-CM2 LA basin
    and produce 4 weighted-mean projections under different policy assumptions.
    """
    from rcmes_mcp.tools.analysis import combine_scenarios_weighted

    scenarios = ["ssp126", "ssp245", "ssp370", "ssp585"]
    scenario_ids = {}
    raw_series = {}  # for plotting individual SSPs underneath
    for sc in scenarios:
        try:
            data = _load("tas", "ACCESS-CM2", sc,
                         "2050-01-01", "2099-12-31", **LA_BASIN)
            ds_id = session_manager.store(
                data=data, source="NEX-GDDP-CMIP6", variable="tas",
                model="ACCESS-CM2", scenario=sc,
                description=f"ACCESS-CM2 {sc} LA basin 2050-2099 (full 3D)",
            )
            scenario_ids[sc] = ds_id
            raw_series[sc] = data.mean(dim=["lat", "lon"]).resample(time="1YS").mean().compute()
        except Exception as e:
            print(f"[deck]   skip {sc}: {e}")

    if len(scenario_ids) < 2:
        raise RuntimeError("need ≥2 scenarios")

    # Define a few policy weight profiles
    policy_weights = {
        "Equal":         {"ssp126": 0.25, "ssp245": 0.25, "ssp370": 0.25, "ssp585": 0.25},
        "Mitigation-aligned (Paris)": {"ssp126": 0.45, "ssp245": 0.35, "ssp370": 0.15, "ssp585": 0.05},
        "Current-policies":           {"ssp126": 0.05, "ssp245": 0.30, "ssp370": 0.40, "ssp585": 0.25},
        "High-emissions risk":        {"ssp126": 0.05, "ssp245": 0.15, "ssp370": 0.30, "ssp585": 0.50},
    }
    # Build weighted projections via the new tool
    policy_series = {}
    for label, w_full in policy_weights.items():
        # Restrict weights to only the scenarios we successfully loaded
        w = {sc: w_full[sc] for sc in scenario_ids if sc in w_full}
        r = combine_scenarios_weighted(scenario_dataset_ids=scenario_ids, weights=w)
        if not r.get("success"):
            print(f"[deck]   combine '{label}' failed: {r.get('error')}")
            continue
        ws = session_manager.get(r["dataset_id"])
        var = ws.name if isinstance(ws, xr.DataArray) else list(ws.data_vars)[0]
        wda = ws if isinstance(ws, xr.DataArray) else ws[var]
        policy_series[label] = wda.mean(dim=["lat", "lon"]).resample(time="1YS").mean().compute()
        print(f"[deck]   policy '{label}' weights={r['weights_normalized']}")

    # Plot
    fig, ax = plt.subplots(figsize=(11, 5.4))
    ssp_colors = {"ssp126": "#1a9850", "ssp245": "#fdae61", "ssp370": "#d73027", "ssp585": "#7f0000"}
    # individual SSPs as faint dashed lines (anomaly vs 2050)
    ref0 = None
    for sc, ts in raw_series.items():
        vals = ts.values
        if float(vals.mean()) > 200:
            vals = vals - 273.15
        if ref0 is None:
            ref0 = float(vals[0])
        ax.plot(ts.time.values, vals - ref0, color=ssp_colors.get(sc, "#888"),
                linewidth=1.0, alpha=0.4, linestyle="--", label=f"{sc} (raw)")
    # policy mixes as bold lines with smoothing
    policy_colors = ["#222", "#1f77b4", "#9467bd", "#d62728"]
    for (label, ts), color in zip(policy_series.items(), policy_colors):
        vals = ts.values
        if float(vals.mean()) > 200:
            vals = vals - 273.15
        anom = vals - ref0
        kernel = np.ones(5) / 5
        smoothed = np.convolve(anom, kernel, mode="same")
        ax.plot(ts.time.values, smoothed, color=color, linewidth=2.5, label=label)
    ax.set_xlabel("Year")
    ax.set_ylabel("Δ T vs 2050 (°C)")
    ax.set_title("Probability-weighted scenario combinations — ACCESS-CM2 LA basin tas\n"
                 "Dashed = individual SSPs; solid = weighted mixes for different policy priors")
    ax.legend(loc="upper left", fontsize=8, ncol=2)
    ax.grid(alpha=0.3)
    return _save(fig, "fig_scenario_weighted_combine.png")


def _real_per_scenario_weighting():
    """
    Per-scenario model weighting: same 3-model ensemble (ACCESS-CM2, GFDL-ESM4,
    MRI-ESM2-0) but show how skill-weighted projections diverge from equal across
    each SSP. Uses historically-trained weights (data already cached).
    """
    from rcmes_mcp.tools.data_access import load_reference_dataset
    from rcmes_mcp.tools.analysis import (
        calculate_model_weights,
        apply_ensemble_weights,
    )

    # 1) Reference + 3-model historical ensemble (likely already cached)
    ref_r = load_reference_dataset(
        source="ncep-reanalysis-monthly",
        start_date="1980-01-01", end_date="2014-12-31",
        lat_min=32.0, lat_max=36.0, lon_min=-120.0, lon_max=-116.0,
    )
    if not ref_r.get("success"):
        raise RuntimeError(f"reference load failed: {ref_r.get('error')}")
    reference_id = ref_r["dataset_id"]

    models = ["ACCESS-CM2", "GFDL-ESM4", "MRI-ESM2-0"]
    series_hist = []
    used = []
    for m in models:
        try:
            data = _load("tas", m, "historical",
                         "1980-01-01", "2014-12-31", **LA_BASIN)
            monthly = data.resample(time="MS").mean()
            try:
                monthly = monthly.convert_calendar("standard", align_on="date", missing=np.nan)
            except Exception:
                pass
            monthly = monthly.assign_coords(model=m).expand_dims("model")
            series_hist.append(monthly)
            used.append(m)
        except Exception as e:
            print(f"[deck]   skip historical {m}: {e}")
    if len(series_hist) < 2:
        raise RuntimeError("need ≥2 historical models")
    hist_ensemble = xr.concat(series_hist, dim="model", join="outer")
    hist_id = session_manager.store(
        data=hist_ensemble, source="NEX-GDDP-CMIP6", variable="tas",
        model=",".join(used), scenario="historical",
        description=f"Hist ensemble {used} 1980-2014 LA basin",
    )

    # 2) Train weights once on historical (use 'combined')
    wres = calculate_model_weights(
        ensemble_dataset_id=hist_id, method="combined",
        reference_dataset_id=reference_id,
        train_start="1980-01-01", train_end="2014-12-31",
        sigma_d=0.5, sigma_s=0.5,
    )
    if not wres.get("success"):
        raise RuntimeError(f"weight training failed: {wres.get('error')}")
    weights = wres["weights"]
    print(f"[deck]   trained weights {used}: {weights}")

    # 3) For each scenario, build the same 3-model ensemble and apply weights
    scenarios = ["ssp126", "ssp245", "ssp585"]
    fig, axes = plt.subplots(1, len(scenarios), figsize=(15, 4.5), sharey=True)
    if len(scenarios) == 1:
        axes = [axes]
    for ax, sc in zip(axes, scenarios):
        sc_series = []
        sc_used = []
        for m in used:
            try:
                data = _load("tas", m, sc, "2050-01-01", "2099-12-31", **LA_BASIN)
                monthly = data.resample(time="MS").mean()
                try:
                    monthly = monthly.convert_calendar("standard", align_on="date", missing=np.nan)
                except Exception:
                    pass
                monthly = monthly.assign_coords(model=m).expand_dims("model")
                sc_series.append(monthly)
                sc_used.append(m)
            except Exception as e:
                print(f"[deck]   skip {m}/{sc}: {e}")
        if len(sc_series) < 2:
            ax.set_title(f"{sc}\n(insufficient data)")
            continue
        sc_ensemble = xr.concat(sc_series, dim="model", join="outer")
        sc_id = session_manager.store(
            data=sc_ensemble, source="NEX-GDDP-CMIP6", variable="tas",
            model=",".join(sc_used), scenario=sc,
            description=f"{sc} ensemble {sc_used} LA basin",
        )

        # Sub-select weights to the models actually present in this scenario
        keep_idx = [used.index(m) for m in sc_used if m in used]
        w_sub = np.array([weights[i] for i in keep_idx], dtype=float)
        w_sub = w_sub / w_sub.sum()

        # Equal projection
        equal_da = sc_ensemble.mean(dim="model")
        eq_ts = equal_da.mean(dim=["lat", "lon"]).resample(time="1YS").mean().compute()

        # Weighted projection via tool
        ar = apply_ensemble_weights(ensemble_dataset_id=sc_id, weights=list(w_sub))
        if not ar.get("success"):
            ax.set_title(f"{sc} (weighting failed)")
            continue
        wds = session_manager.get(ar["dataset_id"])
        wda = wds if isinstance(wds, xr.DataArray) else wds[list(wds.data_vars)[0]]
        w_ts = wda.mean(dim=["lat", "lon"]).resample(time="1YS").mean().compute()

        for ts, color, label in [(eq_ts, "#666", "equal"), (w_ts, "#c4321c", "skill+independence")]:
            vals = ts.values
            if float(vals.mean()) > 200:
                vals = vals - 273.15
            ref0 = float(vals[0])
            anom = vals - ref0
            kernel = np.ones(5) / 5
            smoothed = np.convolve(anom, kernel, mode="same")
            ax.plot(ts.time.values, anom, color=color, linewidth=0.6, alpha=0.3)
            ax.plot(ts.time.values, smoothed, color=color, linewidth=2.0, label=label)
        ax.set_title(sc.upper())
        ax.set_xlabel("Year")
        ax.grid(alpha=0.3)
        ax.legend(loc="upper left", fontsize=8)
    axes[0].set_ylabel("Δ T vs 2050 (°C)")
    fig.suptitle(
        f"Per-scenario projection: equal vs combined weighting "
        f"({len(used)} models, weights trained on historical 1980–2014 vs NCEP)",
        fontsize=11,
    )
    fig.tight_layout()
    return _save(fig, "fig_per_scenario_weighting.png")


def _synth_scenario_weighted_combine():
    """Synthetic placeholder for the scenario combiner figure."""
    years = np.arange(2050, 2100)
    base = np.linspace(0, 1, len(years))
    fig, ax = plt.subplots(figsize=(11, 5.4))
    ssp = {"ssp126": (1.0, "#1a9850"), "ssp245": (2.5, "#fdae61"),
           "ssp370": (3.8, "#d73027"), "ssp585": (4.6, "#7f0000")}
    for sc, (slope, c) in ssp.items():
        ax.plot(years, slope * base, color=c, linestyle="--", linewidth=1.0, alpha=0.4, label=f"{sc} (raw)")
    profiles = {"Equal": [0.25] * 4,
                "Mitigation-aligned (Paris)": [0.45, 0.35, 0.15, 0.05],
                "Current-policies": [0.05, 0.30, 0.40, 0.25],
                "High-emissions risk": [0.05, 0.15, 0.30, 0.50]}
    pcolors = ["#222", "#1f77b4", "#9467bd", "#d62728"]
    slopes = np.array([1.0, 2.5, 3.8, 4.6])
    for (label, w), c in zip(profiles.items(), pcolors):
        eff = float(np.dot(slopes, w))
        ax.plot(years, eff * base, color=c, linewidth=2.5, label=label)
    ax.set_xlabel("Year"); ax.set_ylabel("Δ T vs 2050 (°C)")
    ax.set_title("Probability-weighted scenario combinations (synthetic placeholder)")
    ax.legend(loc="upper left", fontsize=8, ncol=2); ax.grid(alpha=0.3)
    return _save(fig, "fig_scenario_weighted_combine.png")


def _synth_per_scenario_weighting():
    """Synthetic placeholder for per-scenario weighting figure."""
    years = np.arange(2050, 2100)
    base = np.linspace(0, 1, len(years))
    scenarios = [("SSP126", 1.0), ("SSP245", 2.5), ("SSP585", 4.5)]
    fig, axes = plt.subplots(1, 3, figsize=(15, 4.5), sharey=True)
    for ax, (sc, slope) in zip(axes, scenarios):
        ax.plot(years, slope * base, color="#666", linewidth=2, label="equal")
        ax.plot(years, 0.92 * slope * base + 0.1, color="#c4321c", linewidth=2, label="skill+independence")
        ax.set_title(sc); ax.set_xlabel("Year"); ax.grid(alpha=0.3); ax.legend(loc="upper left", fontsize=8)
    axes[0].set_ylabel("Δ T vs 2050 (°C)")
    fig.suptitle("Per-scenario projection: equal vs combined weighting (synthetic placeholder)", fontsize=11)
    fig.tight_layout()
    return _save(fig, "fig_per_scenario_weighting.png")


def _synth_weighting_validation():
    """Two-panel synthetic skill scores + projection comparison."""
    methods = list(_WEIGHT_METHOD_COLORS.keys())
    colors = [_WEIGHT_METHOD_COLORS[m] for m in methods]
    train_rmse = [1.10, 0.78, 1.02, 0.74]
    test_rmse  = [1.18, 0.96, 1.05, 0.88]
    test_corr  = [0.62, 0.74, 0.66, 0.79]

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))
    x = np.arange(len(methods))
    bw = 0.35
    axes[0].bar(x - bw / 2, train_rmse, bw, color=colors, alpha=0.55, label="train")
    axes[0].bar(x + bw / 2, test_rmse, bw, color=colors, label="test")
    axes[0].set_xticks(x); axes[0].set_xticklabels(methods, rotation=15)
    axes[0].set_ylabel("Regional-mean RMSE vs reference (K)")
    axes[0].set_title("Skill — lower is better"); axes[0].legend()
    axes[0].grid(alpha=0.3, axis="y")

    axes[1].bar(x, test_corr, color=colors)
    axes[1].set_xticks(x); axes[1].set_xticklabels(methods, rotation=15)
    axes[1].set_ylabel("Pearson r (test)")
    axes[1].set_title("Correlation — higher is better")
    axes[1].set_ylim(0, 1.05); axes[1].grid(alpha=0.3, axis="y")
    fig.suptitle("Ensemble-weighting validation (synthetic placeholder)", fontsize=11)
    fig.tight_layout()
    _save(fig, "fig_weighting_validation.png")

    # Comparison projections
    years = np.arange(2050, 2100)
    base = np.linspace(0, 4.0, len(years))
    fig2, ax = plt.subplots(figsize=(10, 5))
    multipliers = {"equal": 1.0, "skill": 0.92, "independence": 0.97, "combined": 0.88}
    for m, c in zip(methods, colors):
        proj = multipliers[m] * base + 0.18 * _RNG.standard_normal(len(years))
        kernel = np.ones(5) / 5
        smoothed = np.convolve(proj, kernel, mode="same")
        ax.plot(years, proj, color=c, alpha=0.25, linewidth=0.6)
        ax.plot(years, smoothed, color=c, linewidth=2, label=m)
    ax.set_xlabel("Year"); ax.set_ylabel("Δ T vs 2050 (°C)")
    ax.set_title("Weighted SSP5-8.5 projection comparison (synthetic placeholder)")
    ax.grid(alpha=0.3); ax.legend(loc="upper left")
    return _save(fig2, "fig_weighting_comparison.png")


# ----------------------------------------------------------------------
# Synthetic fallback generators (used when real data is unavailable)
# ----------------------------------------------------------------------


def _synth_anomaly_hovmoller():
    """Standardized-anomaly Hovmöller (time × longitude) with detected events boxed."""
    n_t, n_lon = 240, 120
    times = np.arange("2050-01", "2070-01", dtype="datetime64[M]")[:n_t]
    lons = np.linspace(-130, -110, n_lon)  # California-ish
    base = 0.4 * _RNG.standard_normal((n_t, n_lon))

    # Inject coherent extreme events (positive and negative blobs)
    events = [
        (40, 60, 3.4, 18, 30),   # (t0, lon_idx0, peak, dt, dlon)
        (90, 25, -2.8, 12, 25),
        (150, 80, 3.0, 22, 35),
        (200, 40, -2.2, 10, 20),
    ]
    for t0, x0, peak, dt, dx in events:
        for ti in range(max(0, t0), min(n_t, t0 + dt)):
            for xi in range(max(0, x0), min(n_lon, x0 + dx)):
                rt = (ti - t0 - dt / 2) / (dt / 2)
                rx = (xi - x0 - dx / 2) / (dx / 2)
                base[ti, xi] += peak * np.exp(-(rt ** 2 + rx ** 2))

    fig, ax = plt.subplots(figsize=(8.5, 5.4))
    vlim = 3.5
    im = ax.pcolormesh(lons, times, base, cmap="RdBu_r", vmin=-vlim, vmax=vlim, shading="auto")
    # Draw event boxes
    for t0, x0, peak, dt, dx in events:
        x_left, x_right = lons[x0], lons[min(n_lon - 1, x0 + dx - 1)]
        t_top, t_bot = times[t0], times[min(n_t - 1, t0 + dt - 1)]
        ax.plot(
            [x_left, x_right, x_right, x_left, x_left],
            [t_top, t_top, t_bot, t_bot, t_top],
            color="black", linewidth=1.3, linestyle="--",
        )
    ax.set_xlabel("Longitude (°)")
    ax.set_ylabel("Time")
    ax.set_title("Standardized anomaly Hovmöller — extreme events outlined")
    cb = plt.colorbar(im, ax=ax, pad=0.02)
    cb.set_label("z-score (σ)")
    return _save(fig, "fig_anomaly_hovmoller.png")


def _synth_eof():
    """Two-panel EOF: leading spatial mode (dipole) + PC1 time series."""
    lats = np.linspace(-20, 20, 50)
    lons = np.linspace(150, 270, 130)
    LON, LAT = np.meshgrid(lons, lats)
    pattern = -np.cos(np.deg2rad(LAT)) * np.cos(np.deg2rad(2 * (LON - 200)))

    n_t = 360
    months = np.arange("1980-01", "2010-01", dtype="datetime64[M]")[:n_t]
    pc1 = (
        1.5 * np.sin(2 * np.pi * np.arange(n_t) / 60.0)  # 5-yr-ish ENSO-like
        + 0.6 * _RNG.standard_normal(n_t)
    )

    fig = plt.figure(figsize=(11, 4.5))
    ax_map = fig.add_subplot(1, 2, 1)
    vlim = float(np.nanmax(np.abs(pattern)))
    im = ax_map.pcolormesh(lons, lats, pattern, cmap="RdBu_r", vmin=-vlim, vmax=vlim, shading="auto")
    ax_map.set_xlabel("Longitude (°)")
    ax_map.set_ylabel("Latitude (°)")
    ax_map.set_title("EOF1 — leading spatial mode (42% variance)")
    plt.colorbar(im, ax=ax_map, pad=0.02, label="loading")

    ax_ts = fig.add_subplot(1, 2, 2)
    ax_ts.plot(months, pc1, color="#0f4c81", linewidth=1.2)
    ax_ts.axhline(0, color="grey", linewidth=0.6)
    ax_ts.fill_between(months, pc1, 0, where=pc1 > 0, color="#c4321c", alpha=0.3)
    ax_ts.fill_between(months, pc1, 0, where=pc1 < 0, color="#1a73c2", alpha=0.3)
    ax_ts.set_xlabel("Time")
    ax_ts.set_ylabel("PC1 amplitude")
    ax_ts.set_title("Principal component time series")
    ax_ts.grid(alpha=0.3)
    fig.tight_layout()
    return _save(fig, "fig_eof.png")


def _synth_ensemble_spread():
    """Ensemble spread plot: median + 25-75 + 5-95 bands across 8 model trajectories."""
    n_t = 100
    years = np.arange(2000, 2100)
    # Shared signal (warming) + per-model spread that grows over time
    warming = np.linspace(0, 4.5, n_t)
    n_models = 8
    members = []
    for k in range(n_models):
        m_offset = 0.4 * _RNG.standard_normal()
        slope_var = 1 + 0.18 * _RNG.standard_normal()
        noise = 0.25 * _RNG.standard_normal(n_t)
        members.append(m_offset + slope_var * warming + noise)
    members = np.array(members)
    p05, p25, p50, p75, p95 = np.percentile(members, [5, 25, 50, 75, 95], axis=0)

    fig, ax = plt.subplots(figsize=(9, 5))
    ax.fill_between(years, p05, p95, color="#1f77b4", alpha=0.18, label="5–95% across models")
    ax.fill_between(years, p25, p75, color="#1f77b4", alpha=0.3, label="25–75%")
    for m in members:
        ax.plot(years, m, color="grey", linewidth=0.45, alpha=0.5)
    ax.plot(years, p50, color="#1f77b4", linewidth=2.0, label="Median")
    ax.set_xlabel("Year")
    ax.set_ylabel("Δ Temperature (°C)")
    ax.set_title("Multi-model ensemble spread — SSP5-8.5 (8 models)")
    ax.legend(loc="upper left", frameon=True)
    ax.grid(alpha=0.3)
    return _save(fig, "fig_ensemble_spread.png")


def _synth_scenario_fan():
    """Scenario fan chart: SSP1-2.6 / 2-4.5 / 5-8.5 with 5-95% bands."""
    n_t = 100
    years = np.arange(2000, 2100)
    fig, ax = plt.subplots(figsize=(9, 5))

    scenarios = [
        ("ssp126", 1.4, "#1a9850"),
        ("ssp245", 2.7, "#fdae61"),
        ("ssp585", 4.6, "#7f0000"),
    ]
    n_models = 6
    for label, end_warming, color in scenarios:
        signal = np.linspace(0, end_warming, n_t)
        members = np.array([
            (1 + 0.15 * _RNG.standard_normal()) * signal
            + 0.35 * _RNG.standard_normal() * np.linspace(0, 1, n_t)
            + 0.18 * _RNG.standard_normal(n_t)
            for _ in range(n_models)
        ])
        med = np.median(members, axis=0)
        lo = np.percentile(members, 5, axis=0)
        hi = np.percentile(members, 95, axis=0)
        ax.fill_between(years, lo, hi, color=color, alpha=0.18)
        ax.plot(years, med, color=color, linewidth=1.8, label=f"{label} (median, 5–95%)")

    ax.set_xlabel("Year")
    ax.set_ylabel("Δ Temperature vs 2000 (°C)")
    ax.set_title("Scenario fan chart — divergence across SSP pathways")
    ax.legend(loc="upper left", frameon=True)
    ax.grid(alpha=0.3)
    return _save(fig, "fig_scenario_fan.png")


def _synth_emergence():
    """Time-of-emergence map (year when |signal| > σ_threshold * baseline std)."""
    lats = np.linspace(-60, 75, 70)
    lons = np.linspace(-180, 180, 140)
    LON, LAT = np.meshgrid(lons, lats)
    # Tropics emerge first; arctic emerges later but strongly; mid-latitudes mixed
    base = 2030 + 30 * np.abs(LAT) / 75 + 4 * _RNG.standard_normal(LAT.shape)
    base = np.clip(base, 2025, 2095)
    # Add some never-emerged speckle
    mask = _RNG.random(LAT.shape) < 0.04
    base = np.where(mask, np.nan, base)

    fig, ax = plt.subplots(figsize=(10, 5))
    im = ax.pcolormesh(lons, lats, base, cmap="viridis", vmin=2030, vmax=2090, shading="auto")
    ax.set_xlabel("Longitude (°)")
    ax.set_ylabel("Latitude (°)")
    ax.set_title("Time of emergence — year |signal| > 1σ baseline natural variability")
    cb = plt.colorbar(im, ax=ax, pad=0.02)
    cb.set_label("Emergence year")
    return _save(fig, "fig_time_of_emergence.png")


def _synth_change_and_agreement():
    """Two-panel: ensemble-mean change (left) + model agreement on sign (right)."""
    lats = np.linspace(-60, 75, 70)
    lons = np.linspace(-180, 180, 140)
    LON, LAT = np.meshgrid(lons, lats)
    # Change: positive everywhere, larger at high lats (polar amplification)
    change = 1.0 + 4.5 * np.abs(LAT) / 75 + 0.6 * _RNG.standard_normal(LAT.shape)
    # Agreement: high in tropics & arctic, lower over mid-latitudes & oceans
    agree = 0.6 + 0.35 * (1 - np.exp(-((LAT - 30) / 25) ** 2))
    agree += 0.05 * _RNG.standard_normal(LAT.shape)
    agree = np.clip(agree, 0, 1)

    fig, axes = plt.subplots(1, 2, figsize=(12, 4.6))
    vlim = float(np.nanmax(np.abs(change)))
    im1 = axes[0].pcolormesh(lons, lats, change, cmap="RdBu_r", vmin=-vlim, vmax=vlim, shading="auto")
    axes[0].set_title("Ensemble-mean ΔT (°C, late-century vs baseline)")
    axes[0].set_xlabel("Longitude")
    axes[0].set_ylabel("Latitude")
    plt.colorbar(im1, ax=axes[0], pad=0.02, label="°C")

    im2 = axes[1].pcolormesh(lons, lats, agree, cmap="YlGnBu", vmin=0, vmax=1, shading="auto")
    axes[1].set_title("Model agreement on sign (IPCC stippling)")
    axes[1].set_xlabel("Longitude")
    plt.colorbar(im2, ax=axes[1], pad=0.02, label="fraction agreeing")

    fig.tight_layout()
    return _save(fig, "fig_change_and_agreement.png")


# ----------------------------------------------------------------------
# Dispatchers: try real data, fall back to synthetic on any failure
# ----------------------------------------------------------------------

_SKIP_EXISTING = os.environ.get("RCMES_DECK_SKIP_EXISTING", "0") == "1"


def _try_real(real_fn, synth_fn, label, expected_path=None):
    # Fast path: if a fresh PNG already exists, skip the recompute
    if _SKIP_EXISTING and expected_path and Path(expected_path).exists():
        print(f"[deck] === Skipping '{label}' — using existing {expected_path} ===")
        return expected_path
    if _REAL_AVAILABLE:
        try:
            print(f"\n[deck] === Building '{label}' from real RCMES data ===")
            t0 = __import__("time").perf_counter()
            path = real_fn()
            elapsed = round(__import__("time").perf_counter() - t0, 1)
            print(f"[deck]   ✓ {label} done in {elapsed}s → {path}")
            return path
        except Exception as e:
            print(f"[deck]   ✗ {label} real-data path failed: {e}")
            print(f"[deck]   falling back to synthetic for {label}")
    return synth_fn()


def make_anomaly_hovmoller_figure():
    return _try_real(_real_anomaly_hovmoller, _synth_anomaly_hovmoller,
                     "anomaly Hovmöller", str(FIG_DIR / "fig_anomaly_hovmoller.png"))


def make_eof_figure():
    return _try_real(_real_eof, _synth_eof, "EOF", str(FIG_DIR / "fig_eof.png"))


def make_ensemble_spread_figure():
    return _try_real(_real_ensemble_spread, _synth_ensemble_spread,
                     "ensemble spread", str(FIG_DIR / "fig_ensemble_spread.png"))


def make_scenario_fan_figure():
    return _try_real(_real_scenario_fan, _synth_scenario_fan,
                     "scenario fan", str(FIG_DIR / "fig_scenario_fan.png"))


def make_emergence_figure():
    return _try_real(_real_time_of_emergence, _synth_emergence,
                     "time-of-emergence", str(FIG_DIR / "fig_time_of_emergence.png"))


def make_change_and_agreement_figure():
    return _try_real(_real_change_and_agreement, _synth_change_and_agreement,
                     "change + agreement", str(FIG_DIR / "fig_change_and_agreement.png"))


def make_weighting_validation_figure():
    """Builds BOTH fig_weighting_validation.png and fig_weighting_comparison.png in one call."""
    return _try_real(_real_weighting_validation, _synth_weighting_validation,
                     "weighting validation", str(FIG_DIR / "fig_weighting_comparison.png"))


def make_scenario_weighted_figure():
    return _try_real(_real_scenario_weighted_combine, _synth_scenario_weighted_combine,
                     "scenario combiner", str(FIG_DIR / "fig_scenario_weighted_combine.png"))


def make_per_scenario_weighting_figure():
    return _try_real(_real_per_scenario_weighting, _synth_per_scenario_weighting,
                     "per-scenario weighting", str(FIG_DIR / "fig_per_scenario_weighting.png"))


# Generate all figures up front so slide code can reference paths
FIG_ANOMALY = make_anomaly_hovmoller_figure()
FIG_EOF = make_eof_figure()
FIG_ENSEMBLE = make_ensemble_spread_figure()
FIG_SCENARIO = make_scenario_fan_figure()
FIG_TOE = make_emergence_figure()
FIG_AGREEMENT = make_change_and_agreement_figure()
FIG_WEIGHTING_COMPARISON = make_weighting_validation_figure()  # also writes fig_weighting_validation.png as a side effect
FIG_WEIGHTING_VALIDATION = str(FIG_DIR / "fig_weighting_validation.png")
FIG_SCENARIO_WEIGHTED = make_scenario_weighted_figure()
FIG_PER_SCENARIO_WEIGHTING = make_per_scenario_weighting_figure()


# ============================================================
# SLIDE 1: Cover — Cover_No Image layout
# ph[15]=section label, ph[16]=title, ph[17]=subtitle, ph[18]=author, ph[20]=footer
# ============================================================
slide = prs.slides.add_slide(COVER_NO_IMAGE)
set_placeholder_text(slide, 15, "Project Update")
set_placeholder_text(slide, 16, "RCMES-MCP")
set_placeholder_text(slide, 17, "Enabling AI-Driven Climate Analysis\nwith NASA's NEX-GDDP-CMIP6 Data")
set_placeholder_text(slide, 18, "March 2026")
set_placeholder_text(slide, 20, FOOTER_TEXT)


# ============================================================
# SLIDE 2: What is RCMES-MCP? — Title and Content
# ph[17]=section, ph[0]=title, ph[19]=content, ph[21]=footer
# ============================================================
slide = prs.slides.add_slide(TITLE_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "What is RCMES-MCP?")
set_footer(slide)

ph = None
for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        ph = p
        break
tf = ph.text_frame
add_bullet_content(tf, [
    "An MCP (Model Context Protocol) server that exposes NASA's climate analysis tools to AI agents",
    "Provides direct access to the NEX-GDDP-CMIP6 dataset — 38 TB of downscaled global climate projections at 0.25° resolution",
    "Enables conversational climate analysis: users ask questions in natural language, AI executes the science",
    "Built on NASA JPL's Regional Climate Model Evaluation System (RCMES) heritage",
    "No data downloads, no coding required — AI handles data access, processing, analysis, and visualization",
], font_size=Pt(14))


# ============================================================
# SLIDE 3: The Data — Two Content layout
# ph[17]=section, ph[0]=title, ph[14]=subtitle, ph[19]=left content, ph[20]=right content
# ============================================================
slide = prs.slides.add_slide(TWO_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "The Data: NEX-GDDP-CMIP6")
set_placeholder_text(slide, 14, "38 TB of bias-corrected, downscaled CMIP6 climate projections hosted on AWS Open Data")
set_footer(slide)

# Left content - dataset details
for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        add_bullet_content(tf, [
            "0.25° resolution (~25 km), daily timestep",
            "Global coverage, 1950–2100",
            "35 CMIP6 climate models",
            "Historical + 4 SSP scenarios",
            "Accessed via S3 in Zarr format — no downloads",
        ], font_size=Pt(13))

# Right content - variables
for p in slide.placeholders:
    if p.placeholder_format.idx == 20:
        tf = p.text_frame
        items = [
            "tas / tasmax / tasmin — Temperature",
            "pr — Precipitation",
            "hurs / huss — Humidity",
            "sfcWind — Surface wind speed",
            "rsds / rlds — Radiation (SW & LW)",
        ]
        # Add header
        p0 = tf.paragraphs[0]
        p0.text = "9 Climate Variables:"
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.font.name = FONT
        p0.space_after = Pt(4)
        for item in items:
            pp = tf.add_paragraph()
            pp.text = item
            pp.font.size = Pt(12)
            pp.font.color.rgb = DARK_GRAY
            pp.font.name = FONT
            pp.space_after = Pt(3)


# ============================================================
# SLIDE 4: Capabilities — Title and Content
# ============================================================
slide = prs.slides.add_slide(TITLE_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "47+ Climate Analysis Tools")
set_footer(slide)

for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        items = [
            "Data Access — Load climate data by model, variable, scenario; region/country selection; multi-model ensemble loading",
            "Processing — Temporal & spatial subsetting, regridding, unit conversion, anomaly & standardized-anomaly, country masking",
            "Statistical Analysis — Climatology, trend with significance, bias, RMSE, correlation, EOF/PCA decomposition",
            "Spatiotemporal Detection — 3D connected-component extreme-event cataloguing, time of emergence",
            "Ensemble & Scenario — Multi-model statistics with IPCC-style agreement maps, scenario-comparison fan charts",
            "Ensemble Weighting — Skill / independence / combined weighting (Knutti 2017) validated against NCEP reanalysis observations",
            "Extreme Indices — 22 ETCCDI indices, heatwave & drought analysis",
            "Visualization — Spatial maps, time series, Hovmöller diagrams, fan charts, ensemble spread plots, Taylor diagrams",
        ]
        add_bullet_content(tf, items, font_size=Pt(13))


# ============================================================
# SLIDE 5: How It Works — Title and Content
# ============================================================
slide = prs.slides.add_slide(TITLE_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "How It Works")
set_footer(slide)

for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        items = [
            "User asks a climate question in natural language (e.g., \"What is the heatwave trend in California?\")",
            "AI agent selects and chains appropriate tools via the MCP protocol",
            "Tools access NEX-GDDP-CMIP6 data directly from AWS S3 — no downloads needed",
            "Analysis engine (xarray, Dask, xclim, SciPy) processes the data",
            "Results returned as statistics, maps, time series plots, or derived datasets",
            "Chainable operations: each tool returns a dataset_id for multi-step workflows",
        ]
        add_bullet_content(tf, items, font_size=Pt(13))

# Add flow diagram below using textboxes
# User → MCP Server → AWS S3 → Analysis → Results
flow_labels = ["User / AI Agent", "MCP Server", "AWS S3 Data", "Analysis Engine", "Results"]
flow_colors = [
    RGBColor(0x00, 0x7D, 0xBA),  # blue
    NAVY,
    RGBColor(0x23, 0x72, 0x3B),  # green
    RGBColor(0xE6, 0x7E, 0x22),  # orange
    RGBColor(0xC0, 0x39, 0x2B),  # red
]
box_w = Inches(1.65)
box_h = Inches(0.45)
gap = Inches(0.18)
total = len(flow_labels) * box_w + (len(flow_labels) - 1) * gap
start_x = (prs.slide_width - total) // 2
y = Inches(4.55)

for i, (label, color) in enumerate(zip(flow_labels, flow_colors)):
    x = start_x + i * (box_w + gap)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    stf = shape.text_frame
    stf.vertical_anchor = MSO_ANCHOR.MIDDLE
    sp = stf.paragraphs[0]
    sp.text = label
    sp.font.size = Pt(9)
    sp.font.bold = True
    sp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sp.font.name = FONT
    sp.alignment = PP_ALIGN.CENTER

    if i < len(flow_labels) - 1:
        ax = x + box_w + Emu(10000)
        ay = y + box_h // 2 - Inches(0.08)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, gap - Emu(20000), Inches(0.16))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
        arrow.line.fill.background()


# ============================================================
# SLIDE 6: Example Use Case — Title, Subtitle and Content
# ============================================================
slide = prs.slides.add_slide(TITLE_SUB_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "Example: Heatwave Trend Analysis")
set_placeholder_text(slide, 14, "\"What is the heatwave trend in California under the high-emissions scenario?\"")
set_footer(slide)

for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        steps = [
            "1. AI loads data: load_climate_data(tasmax, ACCESS-CM2, ssp585, California, 2015–2100)",
            "2. AI analyzes heatwaves: analyze_heatwaves(dataset_id, threshold_percentile=90)",
            "3. AI computes trend: calculate_trend(heatwave_id) → significant warming trend with p-value",
            "4. AI visualizes: generate_timeseries_plot(dataset_id, show_trend=True) → plot returned to user",
            "",
            "All steps executed automatically — the user only asks the question",
        ]
        add_bullet_content(tf, steps, font_size=Pt(13))
        # Make last item italic
        last_p = tf.paragraphs[-1]
        last_p.font.italic = True
        last_p.font.color.rgb = RGBColor(0x00, 0x7D, 0xBA)


# ============================================================
# Helper for figure-backed analytics slides (TITLE_ONLY layout)
# Layout: figure on the left half, bullet text on the right.
# ============================================================
def add_figure_slide(title_text, image_path, bullet_items,
                     img_left=Inches(0.35), img_top=Inches(1.55),
                     img_width=Inches(7.6),
                     text_left=Inches(8.2), text_top=Inches(1.55),
                     text_width=Inches(4.85), text_height=Inches(5.4)):
    s = prs.slides.add_slide(TITLE_ONLY)
    set_placeholder_text(s, 17, "RCMES-MCP")
    set_placeholder_text(s, 0, title_text)
    set_footer(s)
    s.shapes.add_picture(image_path, img_left, img_top, width=img_width)

    txt_box = s.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = txt_box.text_frame
    tf.word_wrap = True
    add_bullet_content(tf, bullet_items, font_size=Pt(13), spacing=Pt(8))
    return s


# ============================================================
# SLIDE 7: Advanced Spatiotemporal Analysis — overview
# ============================================================
slide = prs.slides.add_slide(TITLE_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "Advanced Spatiotemporal Analysis")
set_footer(slide)

for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        items = [
            "Pillar 1 — Anomaly detection & dominant patterns: standardized z-score anomalies, "
            "3D connected-component extreme-event cataloguing, EOF/PCA decomposition, Hovmöller visualization",
            "Pillar 2 — Multi-model ensembles & scenario comparison: parallel multi-model loading, "
            "ensemble statistics with IPCC-style agreement maps, pairwise scenario differences, "
            "fan charts and ensemble spread plots, time-of-emergence mapping",
            "10 new tools added to the agent toolkit — chainable through the same dataset_id pattern",
            "Surfaces standard climate-research diagnostics that previously required custom scripting",
        ]
        add_bullet_content(tf, items, font_size=Pt(13))


# ============================================================
# SLIDE 8: Spatiotemporal Extreme-Event Detection
# ============================================================
add_figure_slide(
    "Spatiotemporal Extreme-Event Detection",
    FIG_ANOMALY,
    [
        "calculate_standardized_anomaly — z-scores per day-of-year (or month/season) "
        "vs a baseline window, comparable across regions and seasons",
        "detect_extreme_events — 3D connected-component labeling on |z| > N·σ; "
        "returns an event catalogue (start/end, peak σ, location, footprint, duration)",
        "Use case: heatwave / cold-snap / wet-spell / drought hunting under SSP scenarios",
        "generate_hovmoller — time × lat (or × lon) view that reveals propagating "
        "anomalies and persistent regimes at a glance",
    ],
)


# ============================================================
# SLIDE 9: Dominant Variability Patterns (EOF / PCA)
# ============================================================
add_figure_slide(
    "Dominant Variability Patterns (EOF / PCA)",
    FIG_EOF,
    [
        "calculate_eof — cosine-latitude-weighted SVD that returns the leading "
        "spatial modes plus their principal-component time series",
        "Surfaces dominant patterns of variability "
        "(teleconnection-like dipoles, seasonal cycles, regime shifts)",
        "Two chained dataset_ids: spatial modes (mode, lat, lon) and PCs (time, mode) — "
        "feed straight into generate_map and generate_timeseries_plot",
        "Reports variance explained per mode and cumulative variance",
    ],
    img_width=Inches(7.6),
)


# ============================================================
# SLIDE 10: Multi-Model Ensemble Loading & Spread
# ============================================================
add_figure_slide(
    "Multi-Model Ensembles — Quantifying Uncertainty",
    FIG_ENSEMBLE,
    [
        "load_multi_model_ensemble — parallel-loads up to 10 CMIP6 models for the "
        "same variable/scenario/region into one dataset with a model dimension",
        "generate_ensemble_spread_plot — median + 25–75% inner band + 5–95% outer band "
        "across models over time (with optional per-model overlays)",
        "Spread is the answer to \"how much do models disagree?\" — wide bands signal "
        "low confidence, tight bands signal robust signal",
        "Cap of 10 protects RAM; pre-subset region & time before loading",
    ],
)


# ============================================================
# SLIDE 11: Ensemble Statistics & IPCC-style Agreement
# ============================================================
add_figure_slide(
    "Ensemble Statistics & Model Agreement",
    FIG_AGREEMENT,
    [
        "calculate_ensemble_statistics — reduces ensemble across the model dimension "
        "into ensemble mean, spread (σ), and min/max envelope dataset_ids",
        "With a baseline window, also returns an agreement map: per-cell fraction of "
        "models agreeing with the ensemble-mean SIGN of change — IPCC stippling",
        "Left: ensemble-mean ΔT (late-century vs baseline). Right: model agreement "
        "(0 = full disagreement, 1 = all models agree)",
        "Standard inputs for confidence-aware climate assessments",
    ],
)


# ============================================================
# SLIDE 12: Scenario Comparison
# ============================================================
add_figure_slide(
    "Scenario Comparison Across SSP Pathways",
    FIG_SCENARIO,
    [
        "compare_scenarios — loads the same variable/model/region across multiple "
        "SSP scenarios (max 5) and computes pairwise time-mean differences",
        "generate_scenario_fan_chart — multi-line projection plot with shaded "
        "5–95% bands per scenario; the canonical IPCC-style figure",
        "Reveals when low- vs high-emission pathways diverge — central to mitigation messaging",
        "Smoothing window option for decadal-scale signals; pairs well with "
        "generate_map on each pairwise difference",
    ],
)


# ============================================================
# SLIDE 13: Time of Emergence
# ============================================================
add_figure_slide(
    "Time of Emergence — When Does Climate Change Arrive?",
    FIG_TOE,
    [
        "calculate_time_of_emergence — per grid cell, the first year that the "
        "rolling-mean signal exceeds N·σ of baseline natural variability",
        "Returns a 2D (lat, lon) map of emergence years (NaN where the signal "
        "never emerges within the window)",
        "Tropics typically emerge first (low natural variability); high latitudes "
        "later but with larger absolute change — the colour gradient tells the story",
        "Useful for impact-assessment timelines and adaptation planning",
    ],
)


# ============================================================
# SLIDE 14: Ensemble Weighting — Methods Overview
# ============================================================
slide = prs.slides.add_slide(TITLE_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "Ensemble Weighting — Beyond Equal Votes")
set_footer(slide)

for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        items = [
            "Default IPCC ensemble averaging treats every model as equally credible — "
            "but CMIP6 has lots of shared lineage (e.g., CESM2 ↔ NorESM2 share atmosphere)",
            "Four schemes implemented (Knutti et al. 2017; Brunner & Knutti 2020):  "
            "(1) Equal — democratic baseline;  "
            "(2) Skill — w ∝ exp(−RMSE² / σ²) against reference observations;  "
            "(3) Independence — penalise pairs of models with low inter-model RMSE;  "
            "(4) Combined — skill × independence, AR6 operational style",
            "load_reference_dataset → calculate_model_weights → apply_ensemble_weights → "
            "weighted projections; validate_ensemble_weighting splits time into train/test "
            "to score generalization against held-out observations",
            "Active collaboration with Elias Massound and Arielle on ensemble calibration paper",
        ]
        add_bullet_content(tf, items, font_size=Pt(13))


# ============================================================
# SLIDE 15: Weighting Validation — Skill Scores
# ============================================================
add_figure_slide(
    "Validating Weighting Schemes Against NCEP Reanalysis",
    FIG_WEIGHTING_VALIDATION,
    [
        "validate_ensemble_weighting trains each scheme on 1980–2002 historical and "
        "scores it on the held-out 2003–2014 window against NCEP/NCAR Reanalysis 1",
        "Left: regional-mean RMSE (lower = better); right: Pearson correlation "
        "of monthly anomalies (higher = better)",
        "Skill-aware schemes typically beat the equal vote when one model substantially "
        "outperforms peers on the regional climatology",
        "Honest framing: only 3 models in this demo — the methodology generalizes to the "
        "full 35-model NEX-GDDP-CMIP6 set",
    ],
)


# ============================================================
# SLIDE 16: Weighted SSP5-8.5 Projections — Method Comparison
# ============================================================
add_figure_slide(
    "Weighted Future Projections by Method",
    FIG_WEIGHTING_COMPARISON,
    [
        "Each colour applies a different trained weight set to the SAME SSP5-8.5 "
        "ensemble (3 models, LA basin, 2050–2099) — the spread between curves is "
        "the structural uncertainty introduced by your weighting choice",
        "When weighting methods disagree by ~0.5 °C late-century, that's the "
        "ensemble-weighting uncertainty band you'd quote alongside model spread",
        "Skill and combined typically tilt the projection toward the model that "
        "best reproduces observations; independence damps clustered-model influence",
        "All four projections are saved as new dataset_ids — feed any into "
        "generate_map / calculate_trend / calculate_time_of_emergence as usual",
    ],
)


# ============================================================
# SLIDE 17: Probabilistic Scenario Combining
# ============================================================
add_figure_slide(
    "Probabilistic Scenario Combining",
    FIG_SCENARIO_WEIGHTED,
    [
        "combine_scenarios_weighted — assign prior probabilities to each SSP and "
        "produce a single weighted-mean projection. Use when the question is "
        "\"what's the expected warming under our best guess of which pathway "
        "society will follow?\"",
        "Dashed lines: individual SSPs (raw). Solid lines: 4 different policy "
        "weight profiles applied to the same 4-scenario set",
        "Mitigation-aligned (Paris) puts most weight on ssp126/245; "
        "Current-policies favours ssp245/370; High-emissions risk loads ssp585",
        "Output is a normal dataset_id — chain into generate_map / generate_timeseries_plot "
        "/ detect_extreme_events to do downstream analysis on the weighted projection",
    ],
)


# ============================================================
# SLIDE 18: Per-Scenario Model Weighting
# ============================================================
add_figure_slide(
    "Model Weighting Across Scenarios",
    FIG_PER_SCENARIO_WEIGHTING,
    [
        "Same 3-model ensemble + same historically-trained weights, applied to "
        "each SSP scenario independently — does weighting shift each scenario's "
        "projection in the same way?",
        "If the skill-favoured model warms faster than the unweighted mean, "
        "weighted projections are higher across ALL scenarios (consistent tilt)",
        "Useful diagnostic: lets you carry one trained weight vector through every "
        "downstream analysis — no need to re-train per scenario",
        "Pair with combine_scenarios_weighted to first weight models within each "
        "scenario, then weight scenarios by policy probability",
    ],
)


# ============================================================
# SLIDE 19: Current Status — Two Content
# (was SLIDE 7 — renumbered)
# ============================================================
slide = prs.slides.add_slide(TWO_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "Current Status")
set_placeholder_text(slide, 14, "")
set_footer(slide)

# Left - accomplishments
for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        p0 = tf.paragraphs[0]
        p0.text = "Accomplishments"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.font.name = FONT
        for item in [
            "32+ MCP tools: full climate analysis pipeline",
            "Operational 38 TB data access via AWS S3",
            "22 ETCCDI extreme climate indices",
            "Interactive web UI (React + FastAPI)",
            "Multiple deployment modes (MCP, Web, API, Docker)",
            "Session management with chainable operations",
            "Country-level masking & geographic selection",
        ]:
            pp = tf.add_paragraph()
            pp.text = item
            pp.font.size = Pt(11)
            pp.font.color.rgb = DARK_GRAY
            pp.font.name = FONT
            pp.space_after = Pt(3)

# Right - tech stack
for p in slide.placeholders:
    if p.placeholder_format.idx == 20:
        tf = p.text_frame
        p0 = tf.paragraphs[0]
        p0.text = "Technology Stack"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.font.name = FONT
        stack = [
            ("Core:", "Python 3.10+, xarray, Dask, zarr"),
            ("Climate:", "xclim, SciPy, NetCDF4"),
            ("Geospatial:", "Shapely, GeoPandas, Cartopy"),
            ("Server:", "FastMCP, FastAPI, Uvicorn"),
            ("Frontend:", "React, TypeScript"),
            ("Data:", "AWS S3, NEX-GDDP-CMIP6"),
            ("AI:", "MCP Protocol, Azure OpenAI"),
        ]
        for label, tech in stack:
            pp = tf.add_paragraph()
            run1 = pp.add_run()
            run1.text = f"{label} "
            run1.font.size = Pt(11)
            run1.font.bold = True
            run1.font.color.rgb = RGBColor(0x00, 0x7D, 0xBA)
            run1.font.name = FONT
            run2 = pp.add_run()
            run2.text = tech
            run2.font.size = Pt(11)
            run2.font.color.rgb = DARK_GRAY
            run2.font.name = FONT
            pp.space_after = Pt(3)


# ============================================================
# SLIDE 20: Path Forward — Two Content
# ============================================================
slide = prs.slides.add_slide(TWO_CONTENT)
set_placeholder_text(slide, 17, "RCMES-MCP")
set_placeholder_text(slide, 0, "Path Forward & Continued Support")
set_placeholder_text(slide, 14, "")
set_footer(slide)

# Left - growth opportunities
for p in slide.placeholders:
    if p.placeholder_format.idx == 19:
        tf = p.text_frame
        p0 = tf.paragraphs[0]
        p0.text = "Growth Opportunities"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.font.name = FONT
        for item in [
            "Expand to additional NASA datasets (MERRA-2, ECCO, GLDAS)",
            "Model evaluation with observational datasets",
            "NASA Earthdata integration for authenticated access",
            "Multi-model ensemble analysis & weighting",
            "Community tool for broader climate science use",
            "Reproducible climate assessments for policy support",
        ]:
            pp = tf.add_paragraph()
            pp.text = item
            pp.font.size = Pt(11)
            pp.font.color.rgb = DARK_GRAY
            pp.font.name = FONT
            pp.space_after = Pt(3)

# Right - impact
for p in slide.placeholders:
    if p.placeholder_format.idx == 20:
        tf = p.text_frame
        p0 = tf.paragraphs[0]
        p0.text = "What Continued Support Enables"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.font.name = FONT
        for item in [
            "Democratize access to NASA climate data for non-experts",
            "Accelerate climate research by removing data engineering bottlenecks",
            "Enable rapid AI-assisted climate impact assessments",
            "Lower the barrier to entry across disciplines",
            "Position NASA JPL at the forefront of AI-enabled science",
            "Support decision-makers with actionable climate intelligence",
        ]:
            pp = tf.add_paragraph()
            pp.text = item
            pp.font.size = Pt(11)
            pp.font.color.rgb = DARK_GRAY
            pp.font.name = FONT
            pp.space_after = Pt(3)


# ============================================================
# SLIDE 21: Closing slide
# ============================================================
slide = prs.slides.add_slide(CLOSING)
set_footer(slide)


# ============================================================
# Save
# ============================================================
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
